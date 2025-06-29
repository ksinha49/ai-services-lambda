# ------------------------------------------------------------------------------
# app.py
# ------------------------------------------------------------------------------
"""
Module: app.py
Description:
  Basic Lambda handler for 2-office-extractor.
Version: 1.0.0
"""

from __future__ import annotations
import io
import json
import logging
import os
from typing import Iterable

import boto3
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

__author__ = "Balakrishna"
__version__ = "1.0.0"
__modified_by__ = "Koushik Sinha"

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)
_handler = logging.StreamHandler()
_handler.setFormatter(
    logging.Formatter("%(asctime)s %(levelname)s [%(name)s] %(message)s", "%Y-%m-%dT%H:%M:%S%z")
)
logger.addHandler(_handler)

s3_client = boto3.client("s3")

BUCKET_NAME = os.environ.get("BUCKET_NAME")
OFFICE_PREFIX = os.environ.get("OFFICE_PREFIX", "")
TEXT_DOC_PREFIX = os.environ.get("TEXT_DOC_PREFIX", "text-docs/")

for name in ("OFFICE_PREFIX", "TEXT_DOC_PREFIX"):
    val = globals()[name]
    if val and not val.endswith("/"):
        globals()[name] = val + "/"

if not BUCKET_NAME:
    raise RuntimeError("BUCKET_NAME environment variable must be set")

def _iter_records(event: dict) -> Iterable[dict]:
    for rec in event.get("Records", []):
        yield rec

def _extract_docx(body: bytes) -> list[dict]:
    doc = Document(io.BytesIO(body))
    out: list[dict] = []
    for p in doc.paragraphs:
        if p.text:
            fmt = {"style": getattr(p.style, "name", None)}
            out.append({"text": p.text, "format": fmt})
    return out

def _extract_pptx(body: bytes) -> list[dict]:
    pres = Presentation(io.BytesIO(body))
    slides = []
    for i, slide in enumerate(pres.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text:
                    texts.append(shape.text)
        slides.append({"slide": i, "text": "\n".join(texts)})
    return slides

def _extract_xlsx(body: bytes) -> dict:
    wb = load_workbook(filename=io.BytesIO(body), data_only=True)
    sheets: dict[str, list[list[str]]] = {}
    for sheet in wb:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(["" if cell is None else str(cell) for cell in row])
        sheets[sheet.title] = rows
    return sheets

def _process_record(record: dict) -> None:
    bucket = record.get("s3", {}).get("bucket", {}).get("name")
    key = record.get("s3", {}).get("object", {}).get("key")
    if bucket != BUCKET_NAME or not key:
        logger.info("Skipping record with bucket=%s key=%s", bucket, key)
        return
    if not key.startswith(OFFICE_PREFIX):
        logger.info("Key %s outside prefix %s - skipping", key, OFFICE_PREFIX)
        return

    ext = os.path.splitext(key)[1].lower()
    if ext not in {".docx", ".pptx", ".xlsx"}:
        logger.info("Unsupported extension %s - skipping", ext)
        return

    obj = s3_client.get_object(Bucket=BUCKET_NAME, Key=key)
    body = obj["Body"].read()

    if ext == ".docx":
        content = _extract_docx(body)
        typ = "docx"
    elif ext == ".pptx":
        content = _extract_pptx(body)
        typ = "pptx"
    else:
        content = _extract_xlsx(body)
        typ = "xlsx"

    document_id = os.path.splitext(os.path.basename(key))[0]
    dest_key = f"{TEXT_DOC_PREFIX}{document_id}.json"

    payload = {
        "documentId": document_id,
        "type": typ,
        "content": content,
    }

    s3_client.put_object(
        Bucket=BUCKET_NAME,
        Key=dest_key,
        Body=json.dumps(payload).encode("utf-8"),
        ContentType="application/json",
    )
    logger.info("Wrote %s", dest_key)

def lambda_handler(event: dict, context: dict) -> dict:
    logger.info("Received event for 2-office-extractor: %s", event)
    for rec in _iter_records(event):
        try:
            _process_record(rec)
        except Exception as exc:  # pragma: no cover - runtime safety
            logger.error("Error processing record %s: %s", rec, exc)
    return {
        "statusCode": 200,
        "body": json.dumps({"message": "2-office-extractor executed"})
    }
